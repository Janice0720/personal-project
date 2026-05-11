#!/usr/bin/env python3
"""產生 FinCheck 期中簡報（研究動機 + 四大分析模組）。"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = subtitle


def add_bullet_slide(prs, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
        p.font.name = "PingFang TC"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)


def add_two_column_bullet_slide(prs, title: str, left_title: str, left: list[str], right_title: str, right: list[str]) -> None:
    """使用空白版型手動放兩欄文字（避免依賴特定佈景主題）。"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # 標題
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = tx.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "PingFang TC"

    def fill_box(x, y, w, h, box_title: str, lines: list[str]) -> None:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        btf = box.text_frame
        btf.word_wrap = True
        p0 = btf.paragraphs[0]
        p0.text = box_title
        p0.font.size = Pt(22)
        p0.font.bold = True
        p0.font.name = "PingFang TC"
        for line in lines:
            p = btf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(16)
            p.font.name = "PingFang TC"
            p.space_after = Pt(6)

    fill_box(0.5, 1.2, 4.4, 5.5, left_title, left)
    fill_box(5.1, 1.2, 4.4, 5.5, right_title, right)


def main() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "FinCheck — 個人財務健檢 AI 系統",
        "研究動機與四大分析模組\n在職碩班｜AI × 資料科學 × 金融分析｜期中報告",
    )

    add_bullet_slide(
        prs,
        "為什麼現在談這個？— 報稅季與日常財務的交會",
        [
            "每年約 5 月綜合所得稅結算申報期間，多數人都會整理一整年的「各類所得」：薪資扣繳、股利憑單、利息、租金等。",
            "報稅不只是填表，而是少數會強迫我們把「錢從哪裡來」一次攤開的時刻。",
            "對老師與同學而言：這是具體、有季節感的情境—與課堂上的財稅統計、資料開放、個資授權議題直接扣合。",
            "本專案想回答：在這些數字背後，能否用資料科學與 AI，把「結構」講清楚，並給出可理解的健檢視角？",
        ],
    )

    add_bullet_slide(
        prs,
        "研究動機（一）— 核心問題",
        [
            "題目脈絡：AI × 資料科學 × 金融分析期中報告。",
            "一般民眾難以客觀評估自己的財務健康狀態。",
            "薪資結構是否穩健？所得多元化程度如何？與同層台灣申報戶相比處於什麼分位？",
            "系統定位：使用者透過 MyData 授權下載個人所得資料後，由 AI 系統自動分析所得結構，產出個人化財務健檢報告。",
        ],
    )

    add_bullet_slide(
        prs,
        "研究動機（二）— 痛點拆解",
        [
            "缺乏比較基準：不知道所得結構與同齡或同分位台灣人相比的水準。",
            "所得來源分散：薪資、股利、利息、租賃等分散在不同憑單與帳戶，難以整合看全貌。",
            "節稅盲點：不清楚稅務負擔是否偏高、是否有調整與規劃空間。",
            "一句話痛點：「我的財務健不健康？跟別人比怎麼樣？我該怎麼做？」",
        ],
    )

    add_bullet_slide(
        prs,
        "研究動機（三）— 所得結構與金融資產的連結",
        [
            "個人所得結構可視為「持有金融資產的結果指標」。",
            "股利所得 → 台股／ETF 等權益資產的報酬。",
            "利息所得 → 債券、定存、貨幣型基金等固定收益。",
            "租賃所得 → 不動產資產的租金報酬。",
            "執行業務所得 → 人力資本（Human Capital）作為一種資產。",
            "本系統從跨類別所得結構切入，屬於個人層級的跨資產金融分析視角。",
        ],
    )

    add_two_column_bullet_slide(
        prs,
        "期中報告三面向對應",
        "A 資料工程",
        [
            "財政部 data.gov.tw 公開統計資料",
            "依統計分佈生成合成個體資料（展示用規模）",
            "資料清洗與特徵準備",
        ],
        "B AI／ML 與 C 洞察",
        [
            "B：K-Means 財務人物分群、多維度評分（熵、稅務效率）",
            "C：Streamlit 互動儀表板、個人化健檢報告文字",
        ],
    )

    add_bullet_slide(
        prs,
        "模組 A｜所得多元化評分",
        [
            "方法：Shannon Entropy（夏農熵）衡量所得來源的分散程度。",
            "所得幾乎僅來自薪資 → 熵值偏低 → 評分較低（穩定但來源單一）。",
            "薪資、股利、利息等並重 → 熵值較高 → 評分較高（多元、抗風險結構較佳）。",
            "解讀重點：不是「越多越好」，而是協助使用者看見「集中度 vs. 分散度」。",
        ],
    )

    add_bullet_slide(
        prs,
        "模組 B｜稅務效率評分",
        [
            "指標：稅務負擔率 ＝ 扣繳稅額 ÷ 給付總額（可視情境延伸可扣抵）。",
            "對照財政部公開統計中同分位（或同層）之平均稅務負擔率。",
            "若顯著高於同層平均：可能代表節稅或扣除額運用尚有規劃空間（需個案化解讀）。",
            "與報稅季呼應：民眾在申報時最在意「繳多繳少」與「是否合理」。",
        ],
    )

    add_bullet_slide(
        prs,
        "模組 C｜K-Means 財務人物分群",
        [
            "依所得佔比等特徵，將使用者對應到可解釋的「財務人物」類型，便於溝通與建議。",
            "薪資型：薪資佔比高 → 建議方向：建立被動收入、開始投資。",
            "投資型：股利＋利息佔比高 → 建議方向：資產配置多元化。",
            "自雇型：執行業務所得佔比高 → 建議方向：退休規劃、保障型保險。",
            "混合型：各類所得相對均衡 → 建議方向：持續優化稅務效率與再平衡。",
            "（研究／實作細節可搭配 Elbow、Silhouette 等指標於技術附錄說明。）",
        ],
    )

    add_bullet_slide(
        prs,
        "模組 D｜AI 個人化建議",
        [
            "輸入：健檢分數、分群結果、所得結構摘要。",
            "產出：約 2～3 條個人化建議文字。",
            "實作策略：優先採規則引擎，確保可解釋與穩定；進階可選接 OpenAI API 生成更自然語句。",
            "目標：讓非技術背景的使用者也能帶走「下一步可以思考什麼」。",
        ],
    )

    add_bullet_slide(
        prs,
        "系統流程（摘要）",
        [
            "訓練階段：財政部公開統計 → 合成個體資料 → K-Means 訓練 → 儲存模型。",
            "推論階段：MyData 個人所得（情境模擬亦可）→ 特徵計算 → 套用模型 → 多維評分與報告。",
            "產出：所得結構視覺化、雷達或分項評分、人物類型與 AI 建議。",
        ],
    )

    add_bullet_slide(
        prs,
        "資料來源與合規（簡述）",
        [
            "訓練／對照：財政部財政資訊中心於 data.gov.tw 之綜所稅分位統計、村里所得統計等（政府開放資料授權）。",
            "參考：主計總處家庭收支調查報告（所得結構參數）。",
            "個人端：MyData 個人所得欄位（所得類別、給付總額、扣繳稅額、可扣抵稅額等）；期中展示得以模擬資料替代。",
        ],
    )

    add_title_slide(prs, "謝謝聆聽", "Q & A")

    out = Path(__file__).resolve().parent.parent / "reports" / "FinCheck_研究動機與四大模組_期中簡報.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


if __name__ == "__main__":
    path = main()
    print(path)
